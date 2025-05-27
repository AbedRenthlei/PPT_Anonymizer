import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from pptx.dml.color import RGBColor
import os
try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
except ImportError:
    TkinterDnD = None

def apply_case(original_char, replacement_char):
    """Apply the case of original_char to replacement_char."""
    if original_char.isupper():
        return replacement_char.upper()
    elif original_char.islower():
        return replacement_char.lower()
    return replacement_char

def generate_fixed_text(original_text):
    """Generate fixed text replacement with 'x', preserving spaces, punctuation, and case."""
    fixed_text = "x"
    result = ""
    fixed_idx = 0
    for char in original_text:
        if char.isspace() or not char.isalnum():
            result += char
        else:
            replacement_char = fixed_text[fixed_idx % len(fixed_text)].lower()
            result += apply_case(char, replacement_char)
            fixed_idx += 1
    return result

def copy_formatting(original_run, new_run, shape=None, slide=None):
    """Copy all font formatting from original to new run."""
    if not original_run.font:
        return
    font = original_run.font
    new_font = new_run.font
    # Copy font properties
    new_font.name = font.name
    new_font.size = font.size
    new_font.bold = font.bold
    new_font.italic = font.italic
    new_font.underline = font.underline
    # Copy color
    if font.color and hasattr(font.color, 'type') and font.color.type:
        try:
            if font.color.type == 1:  # RGB
                new_font.color.rgb = font.color.rgb
            elif font.color.type == 2:  # Theme color
                new_font.color.theme_color = font.color.theme_color
                if font.color.brightness:
                    new_font.color.brightness = font.color.brightness
        except (AttributeError, ValueError):
            pass
    # Check run XML for color
    run_element = original_run._r
    if run_element.xpath('.//a:solidFill'):
        fill = run_element.xpath('.//a:solidFill')[0]
        if fill.xpath('.//a:srgbClr'):
            rgb = fill.xpath('.//a:srgbClr')[0].get('val')
            new_font.color.rgb = RGBColor.from_string(rgb)
        elif fill.xpath('.//a:schemeClr'):
            scheme = fill.xpath('.//a:schemeClr')[0].get('val')
            try:
                new_font.color.theme_color = scheme
            except ValueError:
                pass

def replace_text_in_shape(shape, slide=None):
    """Replace text in a shape with 'x'."""
    if shape.shape_type == 6:  # Group shape
        for sub_shape in shape.shapes:
            replace_text_in_shape(sub_shape, slide)
        return
    
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return
    
    # Preserve text box properties
    if shape.text_frame:
        shape.text_frame.auto_size = shape.text_frame.auto_size
        shape.text_frame.word_wrap = shape.text_frame.word_wrap
        shape.text_frame.margin_left = shape.text_frame.margin_left
        shape.text_frame.margin_right = shape.text_frame.margin_right
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            if not isinstance(original_text, str) or not original_text.strip():
                continue
            run.text = generate_fixed_text(original_text)
            copy_formatting(run, run, shape, slide)

def replace_text_in_table(table, slide=None):
    """Replace text in a table with 'x'."""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    if not isinstance(original_text, str) or not original_text.strip():
                        continue
                    run.text = generate_fixed_text(original_text)
                    copy_formatting(run, run, cell, slide)

def anonymize_ppt(input_path, output_path):
    """Anonymize text in a PowerPoint presentation with 'x'."""
    try:
        prs = Presentation(input_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    replace_text_in_shape(shape, slide)
                elif shape.has_table:
                    replace_text_in_table(shape.table, slide)
                elif shape.shape_type == 6:
                    replace_text_in_shape(shape, slide)
        prs.save(output_path)
        return True, f"Successfully anonymized presentation saved to: {output_path}"
    except PermissionError:
        return False, "Error: Output file is open or in use. Close it or choose a different name."
    except Exception as e:
        return False, f"Error processing the presentation: {str(e)}"

class PPTAnonymizerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Text Anonymizer")
        self.root.configure(bg="#F5F6F5")
        self.root.minsize(400, 600)

        # Configure style
        style = ttk.Style()
        style.configure("TLabel", background="#F5F6F5", font=("Helvetica", 12))
        style.configure("TEntry", font=("Helvetica", 11))

        # Main container frame
        container = ttk.Frame(self.root)
        container.pack(fill="both", expand=True)

        # Center main frame
        main_frame = ttk.Frame(container, padding="20")
        main_frame.pack(expand=True, pady=20, padx=20)

        # Variables
        self.input_path = tk.StringVar()
        self.output_path = tk.StringVar()

        # Title
        ttk.Label(main_frame, text="PowerPoint Text Anonymizer", font=("Helvetica", 18, "bold")).pack(pady=(0, 20))

        # Drag-and-Drop Area
        self.drop_frame = tk.Frame(main_frame, bg="#E8ECEF", bd=1, relief="solid", width=300, height=80)
        self.drop_frame.pack_propagate(False)
        self.drop_frame.pack(padx=10, pady=10)
        self.drop_label = tk.Label(self.drop_frame, text="Drag and drop .pptx file here or click to browse", 
                                  bg="#E8ECEF", fg="#333333", font=("Helvetica", 11), wraplength=280)
        self.drop_label.pack(expand=True, fill="both", padx=10, pady=10)
        self.drop_label.bind("<Button-1>", lambda e: self.browse_input())

        if TkinterDnD:
            self.root.drop_target_register(DND_FILES)
            self.root.dnd_bind('<<Drop>>', self.handle_drop)

        # Input File Display
        ttk.Label(main_frame, text="Selected File:").pack(pady=(10, 5))
        self.input_entry = ttk.Entry(main_frame, textvariable=self.input_path, width=30, state="readonly")
        self.input_entry.pack(padx=10)

        # Output File Display
        ttk.Label(main_frame, text="Output File:").pack(pady=(10, 5))
        self.output_entry = ttk.Entry(main_frame, textvariable=self.output_path, width=30, state="readonly")
        self.output_entry.pack(padx=10)

        # Process Button
        self.process_button = tk.Button(main_frame, text="Anonymize", command=self.process, 
                                       bg="#0078D7", fg="white", font=("Helvetica", 12, "bold"), 
                                       relief="flat", padx=20, pady=10)
        self.process_button.pack(pady=20)
        self.process_button.bind("<Enter>", lambda e: self.process_button.config(bg="#005EA6"))
        self.process_button.bind("<Leave>", lambda e: self.process_button.config(bg="#0078D7"))

    def truncate_filename(self, filename, max_length=25):
        """Truncate filename with '...' if too long."""
        if len(filename) > max_length:
            return filename[:max_length-3] + "..."
        return filename

    def handle_drop(self, event):
        file_path = event.data.strip('{}')
        if file_path.endswith('.pptx'):
            self.input_path.set(file_path)
            default_output = os.path.splitext(file_path)[0] + "_anonymized.pptx"
            output_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx")],
                initialfile=os.path.basename(default_output),
                initialdir=os.path.dirname(file_path)
            )
            if output_path:
                self.output_path.set(output_path)
                self.drop_label.config(text="File loaded: " + self.truncate_filename(os.path.basename(file_path)))
            else:
                self.input_path.set("")
                self.drop_label.config(text="Drag and drop .pptx file here or click to browse")
                messagebox.showwarning("Warning", "Output file not selected. Please try again.")
        else:
            messagebox.showerror("Error", "Please drop a valid .pptx file.")

    def browse_input(self):
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
        if file_path:
            self.input_path.set(file_path)
            default_output = os.path.splitext(file_path)[0] + "_anonymized.pptx"
            output_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx")],
                initialfile=os.path.basename(default_output),
                initialdir=os.path.dirname(file_path)
            )
            if output_path:
                self.output_path.set(output_path)
                self.drop_label.config(text="File loaded: " + self.truncate_filename(os.path.basename(file_path)))
            else:
                self.input_path.set("")
                self.drop_label.config(text="Drag and drop .pptx file here or click to browse")
                messagebox.showwarning("Warning", "Output file not selected. Please try again.")

    def process(self):
        input_path = self.input_path.get()
        output_path = self.output_path.get()
        
        if not input_path or not os.path.exists(input_path):
            messagebox.showerror("Error", "Please select a valid input .pptx file.")
            return
        if not output_path:
            messagebox.showerror("Error", "Please select an output file name.")
            return
        if not input_path.endswith(".pptx"):
            messagebox.showerror("Error", "Input file must be a .pptx file.")
            return
        
        success, message = anonymize_ppt(input_path, output_path)
        if success:
            messagebox.showinfo("Success", message)
        else:
            messagebox.showerror("Error", message)

if __name__ == "__main__":
    if TkinterDnD:
        root = TkinterDnD.Tk()
    else:
        root = tk.Tk()
        messagebox.showwarning("Warning", "Drag-and-drop not available. Install tkinterdnd2 for full functionality.")
    app = PPTAnonymizerApp(root)
    root.mainloop()